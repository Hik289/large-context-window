"""
PowerPoint document processor for .ppt and .pptx files.
"""

from pathlib import Path
from typing import List
from agent_memory.core.segment import Segment
from agent_memory.processors.base_processor import FileProcessor, detect_file_type


class PowerPointProcessor(FileProcessor):
    """Processor for Microsoft PowerPoint files (.ppt, .pptx)."""

    def can_process(self, file_path: Path) -> bool:
        """Return ``True`` for PowerPoint files."""
        return detect_file_type(file_path) == "powerpoint"

    def process(self, file_path: Path) -> List[Segment]:
        """Convert a PowerPoint deck into per-title-group segments.

        Args:
            file_path: Path to the PowerPoint file.

        Returns:
            Ordered list of :class:`Segment` objects.
        """
        if not file_path.exists():
            raise FileNotFoundError(f"File not found: {file_path}")

        try:
            from pptx import Presentation
        except ImportError:
            raise ImportError(
                "python-pptx is required to process PowerPoint files. Install with: pip install python-pptx"
            )

        try:
            prs = Presentation(file_path)
            segments: List[Segment] = []
            base_metadata = self._create_base_metadata(file_path, "powerpoint")

            current_title = None
            current_slide_numbers: List[int] = []
            current_contents: List[str] = []

            def flush_current_segment() -> None:
                """Emit one slide-group segment if there is anything pending."""
                nonlocal current_title, current_slide_numbers, current_contents
                if current_title is None:
                    return

                segments.append(
                    Segment(
                        content="\n\n".join(current_contents).strip(),
                        segment_type="slide_group",
                        metadata={
                            **base_metadata,
                            "title": current_title,
                            "slide_numbers": current_slide_numbers.copy(),
                            "slide_number_start": current_slide_numbers[0],
                            "slide_number_end": current_slide_numbers[-1],
                            "slide_count": len(current_slide_numbers),
                        },
                    )
                )

                current_title = None
                current_slide_numbers = []
                current_contents = []

            for slide_num, slide in enumerate(prs.slides, 1):
                slide_texts = [
                    shape.text.strip()
                    for shape in slide.shapes
                    if hasattr(shape, "text") and shape.text.strip()
                ]

                slide_title = slide_texts[0] if slide_texts else f"Slide {slide_num}"
                slide_content = "\n".join(slide_texts) if slide_texts else slide_title

                if current_title is None:
                    # Opening the very first group.
                    current_title = slide_title
                    current_slide_numbers = [slide_num]
                    current_contents = [slide_content]
                elif current_title == slide_title:
                    # Same title as before — extend the active group.
                    current_slide_numbers.append(slide_num)
                    current_contents.append(slide_content)
                else:
                    # Title changed — finalize previous group and start a new one.
                    flush_current_segment()
                    current_title = slide_title
                    current_slide_numbers = [slide_num]
                    current_contents = [slide_content]

            flush_current_segment()

            return segments

        except Exception as exc:
            raise ValueError(f"Failed to process PowerPoint file: {exc}")
